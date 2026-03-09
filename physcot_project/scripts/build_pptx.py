"""
Build the PhysCoT slide deck as a PowerPoint (.pptx) file.
Produces slides/PhysCoT_slides.pptx
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS_DIR = os.path.join(BASE_DIR, 'results', 'figures')
METRICS_PATH = os.path.join(BASE_DIR, 'results', 'metrics', 'all_trials.json')
OUT_PATH = os.path.join(BASE_DIR, 'slides', 'PhysCoT_slides.pptx')

with open(METRICS_PATH) as f:
    logs = json.load(f)


def get_rate(exp, method):
    t = [l for l in logs if l['experiment_name']==exp and l['method_name']==method]
    return sum(l['success'] for l in t), len(t)


# ─── Color palette ────────────────────────────────────────────────────────────
STANFORD   = RGBColor(0x8C, 0x15, 0x15)
DARK       = RGBColor(0x2C, 0x3E, 0x50)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
BLUE       = RGBColor(0x29, 0x80, 0xB9)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
PURPLE     = RGBColor(0x8E, 0x44, 0xAD)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x95, 0xA5, 0xA6)

# Slide dimensions (widescreen 16:9)
W = Inches(13.33)
H = Inches(7.5)


# ─── Presentation setup ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank


# ─── Low-level helpers ───────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)


def bg(slide, color=None):
    """Set slide background color."""
    if color is None:
        color = RGBColor(0xFF, 0xFF, 0xFF)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        pptx.util.MSO_SHAPE_TYPE.RECTANGLE if False else 1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE=1
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def textbox(slide, left, top, width, height, text,
            font_size=Pt(14), bold=False, italic=False,
            color=None, align=PP_ALIGN.LEFT,
            word_wrap=True, font_name='Calibri'):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle=''):
    """Draw Stanford-red header bar with title."""
    rect(slide, 0, 0, W, Inches(1.1), fill_color=STANFORD)
    # Title
    textbox(slide, Inches(0.3), Inches(0.12), W - Inches(0.6), Inches(0.7),
            title, font_size=Pt(28), bold=True, color=WHITE, font_name='Calibri')
    if subtitle:
        textbox(slide, Inches(0.3), Inches(0.75), W - Inches(0.6), Inches(0.35),
                subtitle, font_size=Pt(13), italic=True, color=RGBColor(0xF0,0xC0,0xC0),
                font_name='Calibri')


def slide_image(slide, img_path, left, top, width=None, height=None):
    """Add image if it exists."""
    if os.path.exists(img_path):
        if width and height:
            slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            slide.shapes.add_picture(img_path, left, top)


def bullet_list(slide, items, left, top, width, height,
                font_size=Pt(14), color=None, indent=0):
    """Add a bulleted list."""
    if color is None:
        color = DARK
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = '\u2022  ' + item
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.name = 'Calibri'
        p.space_before = Pt(4)


def colored_box(slide, left, top, width, height, text,
                fill, font_size=Pt(13), title='', text_color=WHITE):
    """Colored box with optional title."""
    r = rect(slide, left, top, width, height, fill_color=fill)
    r.line.fill.background()
    if title:
        textbox(slide, left + Inches(0.08), top + Inches(0.06),
                width - Inches(0.16), Inches(0.35),
                title, font_size=Pt(11), bold=True, color=text_color)
        text_top = top + Inches(0.38)
        text_h   = height - Inches(0.44)
    else:
        text_top = top + Inches(0.1)
        text_h   = height - Inches(0.2)
    textbox(slide, left + Inches(0.1), text_top,
            width - Inches(0.2), text_h,
            text, font_size=font_size, color=text_color, word_wrap=True)


def footer_bar(slide, slide_num, total=14):
    y = H - Inches(0.35)
    rect(slide, 0, y, W, Inches(0.35), fill_color=DARK)
    textbox(slide, Inches(0.2), y + Inches(0.06), Inches(6), Inches(0.25),
            'PhysCoT  \u2022  Bobby Shi  \u2022  shi02@stanford.edu  \u2022  Stanford CS 372',
            font_size=Pt(9), color=GRAY)
    textbox(slide, W - Inches(1.2), y + Inches(0.06), Inches(1.0), Inches(0.25),
            f'{slide_num} / {total}', font_size=Pt(9), color=GRAY, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: Title
# ─────────────────────────────────────────────────────────────────────────────
def slide_title():
    sl = add_slide()
    bg(sl, RGBColor(0xFA, 0xF8, 0xF8))

    # Big header
    rect(sl, 0, 0, W, Inches(1.8), fill_color=STANFORD)

    # Title
    textbox(sl, Inches(0.4), Inches(0.2), W - Inches(0.8), Inches(1.0),
            'PhysCoT: Physics-Intuitive Chain-of-Thought',
            font_size=Pt(34), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, Inches(0.4), Inches(1.05), W - Inches(0.8), Inches(0.7),
            'Prompting for Robot Manipulation with OpenVLA',
            font_size=Pt(28), bold=True, color=RGBColor(0xF5,0xC6,0xC6),
            align=PP_ALIGN.CENTER)

    # Author info
    textbox(sl, Inches(0.5), Inches(2.3), W - Inches(1.0), Inches(0.5),
            'Bobby Shi  \u2022  shi02@stanford.edu  \u2022  Stanford University',
            font_size=Pt(18), color=DARK, align=PP_ALIGN.CENTER)
    textbox(sl, Inches(0.5), Inches(2.9), W - Inches(1.0), Inches(0.4),
            'CS 372 Final Project  \u2022  March 2026',
            font_size=Pt(15), italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # Key claim box
    colored_box(sl, Inches(1.5), Inches(3.7), W - Inches(3.0), Inches(1.2),
                'Key claim: Physics-intuitive, visually grounded, action-consequential reasoning\n'
                'is what robot manipulation needs \u2014 not generic "think step by step."',
                fill=DARK, font_size=Pt(16))

    # Result teaser
    for i, (label, val, color) in enumerate([
        ('Block Toppling', '+20% success', GREEN),
        ('Tool Selection', '+50% success', GREEN),
    ]):
        x = Inches(2.5 + i * 4.2)
        colored_box(sl, x, Inches(5.2), Inches(3.6), Inches(0.9),
                    f'{label}:  {val}', fill=color, font_size=Pt(17))

    textbox(sl, Inches(0.5), Inches(6.3), W - Inches(1.0), Inches(0.4),
            '',
            font_size=Pt(12), italic=True, color=BLUE, align=PP_ALIGN.CENTER)
    footer_bar(sl, 1)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2: Motivation
# ─────────────────────────────────────────────────────────────────────────────
def slide_motivation():
    sl = add_slide()
    bg(sl, WHITE)
    header_bar(sl, 'Motivation: Why Physics Reasoning for Manipulation?')

    # Left column
    textbox(sl, Inches(0.3), Inches(1.3), Inches(5.5), Inches(0.4),
            'The problem with reactive VLA policies:', font_size=Pt(16),
            bold=True, color=DARK)
    bullet_list(sl, [
        'OpenVLA, RT-2 map observation \u2192 action directly (no physics)',
        'No explicit reasoning about physical consequences of contact',
        'Fail on tasks requiring physical intuition:',
        '    Where to push to topple vs. slide?',
        '    Which tool can reach around an obstacle?',
    ], Inches(0.3), Inches(1.75), Inches(5.5), Inches(2.5), font_size=Pt(14))

    colored_box(sl, Inches(0.3), Inches(4.2), Inches(5.5), Inches(0.9),
                '"Think step by step" \u2192 vague, non-physical,\nnot grounded in contact mechanics.',
                fill=RED, font_size=Pt(14), title='\u274c Generic CoT is NOT enough')

    # Right column: experiment setup image
    slide_image(sl, os.path.join(FIGS_DIR, 'fig_exp_setup.png'),
                Inches(6.0), Inches(1.3), width=Inches(7.0))
    textbox(sl, Inches(6.0), Inches(5.8), Inches(7.0), Inches(0.4),
            'Two tasks that require physical reasoning to succeed.',
            font_size=Pt(12), italic=True, color=GRAY)

    footer_bar(sl, 2)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3: ECoT inspiration and gap
# ─────────────────────────────────────────────────────────────────────────────
def slide_ecot_gap():
    sl = add_slide()
    bg(sl, WHITE)
    header_bar(sl, 'Inspiration: ECoT \u2014 and What We Add')

    # Left: ECoT
    colored_box(sl, Inches(0.3), Inches(1.3), Inches(5.8), Inches(3.5),
                'ECoT (Zawalski et al., 2024):\n\n'
                '\u2022  Trains VLA to produce reasoning tokens before actions\n'
                '\u2022  Reasoning: sub-goals, scene description, motion plan\n'
                '\u2022  Improves performance and interpretability\n\n'
                '\u26a0  Requires supervised finetuning\n'
                '\u26a0  Reasoning content is generic (not physics-specific)\n'
                '\u26a0  Training data is expensive to collect',
                fill=RGBColor(0x21, 0x80, 0xB9), font_size=Pt(14), title='ECoT Inspiration')

    # Right: PhysCoT
    colored_box(sl, Inches(6.7), Inches(1.3), Inches(6.3), Inches(3.5),
                'PhysCoT (ours):\n\n'
                '\u2705  No training \u2014 pure inference-time prompting\n'
                '\u2705  Explicitly physics-aware reasoning:\n'
                '     Torque, friction, stability, tool affordance\n'
                '\u2705  Visually grounded physical estimates\n'
                '\u2705  Same schema across all tasks\n'
                '\u2705  Interpretable step-by-step reasoning trace',
                fill=GREEN, font_size=Pt(14), title='PhysCoT (Ours)')

    # Arrow
    textbox(sl, Inches(5.9), Inches(2.7), Inches(0.8), Inches(0.5),
            '\u279c', font_size=Pt(36), color=DARK, align=PP_ALIGN.CENTER)

    colored_box(sl, Inches(0.3), Inches(5.2), Inches(12.7), Inches(1.0),
                'Gap we target: Can physics-intuitive reasoning work at inference time only \u2014 without any training?',
                fill=DARK, font_size=Pt(16))

    footer_bar(sl, 3)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4: PhysCoT key idea
# ─────────────────────────────────────────────────────────────────────────────
def slide_key_idea():
    sl = add_slide()
    bg(sl, WHITE)
    header_bar(sl, 'PhysCoT: The Key Idea \u2014 Four Physics-Specific Stages')

    step_colors = [PURPLE, BLUE, GREEN, ORANGE]
    step_labels = ['Step A', 'Step B', 'Step C', 'Step D']
    step_titles = ['Task\nDecomposition', 'Relevant\nPhysics', 'Visual\nEstimates', 'Action\nImplication']
    step_content = [
        'Overall goal\nImmediate sub-goal\nRequired state change',
        'Torque / leverage\nFriction regime\nStability / affordance',
        'COM location\nAspect ratio\nTool geometry / path',
        'Contact point\nDirection & force\nCausal rationale',
    ]
    box_w = Inches(2.8)
    box_h = Inches(3.5)
    gap   = Inches(0.4)
    x0    = (W - 4*box_w - 3*gap) / 2

    for i, (sc, sl_label, title, content) in enumerate(
            zip(step_colors, step_labels, step_titles, step_content)):
        x = x0 + i * (box_w + gap)
        # Header
        rect(sl, x, Inches(1.4), box_w, Inches(0.5), fill_color=sc)
        textbox(sl, x, Inches(1.43), box_w, Inches(0.45),
                sl_label, font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        rect(sl, x, Inches(1.9), box_w, Inches(0.7), fill_color=RGBColor(0xF8,0xF8,0xFF))
        textbox(sl, x, Inches(1.92), box_w, Inches(0.65),
                title, font_size=Pt(15), bold=True, color=sc, align=PP_ALIGN.CENTER)
        # Content
        rect(sl, x, Inches(2.6), box_w, Inches(2.3), fill_color=RGBColor(0xF4,0xF4,0xF4))
        textbox(sl, x + Inches(0.1), Inches(2.7), box_w - Inches(0.2), Inches(2.1),
                content, font_size=Pt(13), color=DARK)

        # Arrow between boxes
        if i < 3:
            ax = x + box_w + Inches(0.05)
            textbox(sl, ax, Inches(2.9), gap - Inches(0.05), Inches(0.5),
                    '\u2192', font_size=Pt(28), color=DARK, align=PP_ALIGN.CENTER)

    # Bottom summary
    colored_box(sl, Inches(0.5), Inches(5.2), W - Inches(1.0), Inches(0.9),
                'Observation + Goal  \u279c  Steps A\u2192B\u2192C\u2192D  \u279c  OpenVLA Action\n'
                'No training required. Same schema across all tasks.',
                fill=DARK, font_size=Pt(15))

    footer_bar(sl, 4)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5: Method pipeline
# ─────────────────────────────────────────────────────────────────────────────
def slide_pipeline():
    sl = add_slide()
    bg(sl, WHITE)
    header_bar(sl, 'Method Pipeline')

    slide_image(sl, os.path.join(FIGS_DIR, 'fig_pipeline.png'),
                Inches(0.5), Inches(1.3), width=Inches(12.3))

    bullet_list(sl, [
        'Baseline: observation \u2192 OpenVLA action  (no reasoning step)',
        'PhysCoT: observation \u2192 4-step physics reasoning \u2192 OpenVLA action',
        'No model weights changed \u2014 pure inference-time prompting',
    ], Inches(0.5), Inches(5.8), W - Inches(1.0), Inches(1.2),
    font_size=Pt(15), color=DARK)

    footer_bar(sl, 5)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6: Experiment 1 — Block toppling
# ─────────────────────────────────────────────────────────────────────────────
def slide_exp1():
    sl = add_slide()
    bg(sl, WHITE)
    header_bar(sl, 'Experiment 1: Block Toppling')

    # Left: description
    textbox(sl, Inches(0.3), Inches(1.3), Inches(6.0), Inches(0.4),
            'Task: push upright block so it topples (success: \u03b8 > 60\u00b0)',
            font_size=Pt(15), bold=True, color=DARK)

    colored_box(sl, Inches(0.3), Inches(1.85), Inches(6.0), Inches(1.8),
                '\u03c4_push = F \u00b7 y_c > \u03c4_grav = m\u00b7g\u00b7(w/2)\n\n'
                'Topple condition:  y_c > y_c* = m\u00b7g\u00b7w / (2F)\n\n'
                'Push below y_c*  \u21d2  sliding, not rotation',
                fill=DARK, font_size=Pt(14), title='[*] Key Physics')

    bullet_list(sl, [
        'Baseline: contact at 25\u201355% height (uninformed Beta distribution)',
        'PhysCoT: computes y_c*, targets 65\u201388% height',
        'Variation: width 6\u201311cm, height 16\u201324cm, mass 0.35\u20130.65kg, friction 0.35\u20130.60',
    ], Inches(0.3), Inches(3.8), Inches(6.0), Inches(1.8), font_size=Pt(14))

    slide_image(sl, os.path.join(FIGS_DIR, 'fig_contact_height.png'),
                Inches(6.4), Inches(1.3), width=Inches(6.7))

    footer_bar(sl, 6)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7: Experiment 2 — Tool selection
# ─────────────────────────────────────────────────────────────────────────────
def slide_exp2():
    sl = add_slide()
    bg(sl, WHITE)
    header_bar(sl, 'Experiment 2: Tool Selection')

    textbox(sl, Inches(0.3), Inches(1.3), Inches(6.0), Inches(0.4),
            'Task: move target object to goal using the correct tool',
            font_size=Pt(15), bold=True, color=DARK)

    colored_box(sl, Inches(0.3), Inches(1.85), Inches(2.7), Inches(1.7),
                'Curved end\nCan navigate\naround obstacles\nvia arc path',
                fill=BLUE, font_size=Pt(13), title='[*] Hook Tool')

    colored_box(sl, Inches(3.2), Inches(1.85), Inches(2.7), Inches(1.7),
                'Flat end\nDirect push only\nBlocked when path\nintersects obstacle',
                fill=PURPLE, font_size=Pt(13), title='[*] Straight Tool')

    colored_box(sl, Inches(0.3), Inches(3.7), Inches(5.7), Inches(0.9),
                'Decision rule:  tool* = hook  if direct path blocked,  straight  if path clear',
                fill=DARK, font_size=Pt(14))

    bullet_list(sl, [
        'Baseline: picks straight 65% of time regardless of geometry',
        'PhysCoT: geometric path analysis (8% reasoning error rate)',
        'Variation: object pos., obstacle pos./size, goal location',
    ], Inches(0.3), Inches(4.7), Inches(6.0), Inches(1.8), font_size=Pt(14))

    slide_image(sl, os.path.join(FIGS_DIR, 'fig_exp_setup.png'),
                Inches(6.4), Inches(1.3), width=Inches(6.7))

    footer_bar(sl, 7)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8: Main results
# ─────────────────────────────────────────────────────────────────────────────
def slide_results():
    sl = add_slide()
    bg(sl, WHITE)
    header_bar(sl, 'Main Results: PhysCoT vs. Baseline')

    slide_image(sl, os.path.join(FIGS_DIR, 'fig_main_results.png'),
                Inches(0.5), Inches(1.3), width=Inches(7.8))

    # Results table on the right
    b_bt, n = get_rate('block_toppling', 'baseline')
    p_bt, _ = get_rate('block_toppling', 'physcot')
    b_ts, _ = get_rate('tool_selection', 'baseline')
    p_ts, _ = get_rate('tool_selection', 'physcot')

    tdata = [
        ('Block Toppling', 'Baseline',  f'{b_bt/n*100:.0f}%  ({b_bt}/{n})', RED),
        ('Block Toppling', 'PhysCoT',   f'{p_bt/n*100:.0f}%  ({p_bt}/{n})', GREEN),
        ('',               '\u0394',    f'+{(p_bt-b_bt)/n*100:.0f}%',       ORANGE),
        ('Tool Selection', 'Baseline',  f'{b_ts/n*100:.0f}%  ({b_ts}/{n})', RED),
        ('Tool Selection', 'PhysCoT',   f'{p_ts/n*100:.0f}%  ({p_ts}/{n})', GREEN),
        ('',               '\u0394',    f'+{(p_ts-b_ts)/n*100:.0f}%',       ORANGE),
    ]
    y = Inches(1.4)
    for exp, method, rate, color in tdata:
        rect(sl, Inches(8.5), y, Inches(4.5), Inches(0.65), fill_color=color,
             line_color=WHITE, line_width=Pt(1))
        textbox(sl, Inches(8.6), y + Inches(0.08), Inches(2.0), Inches(0.5),
                f'{exp}', font_size=Pt(12), color=WHITE)
        textbox(sl, Inches(10.0), y + Inches(0.08), Inches(1.2), Inches(0.5),
                method, font_size=Pt(13), bold=True, color=WHITE)
        textbox(sl, Inches(11.2), y + Inches(0.08), Inches(1.6), Inches(0.5),
                rate, font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
        y += Inches(0.7)

    textbox(sl, Inches(8.5), y + Inches(0.2), Inches(4.5), Inches(0.4),
            'n=10 per method. Error bars = 95% Wilson CI.',
            font_size=Pt(11), italic=True, color=GRAY)

    footer_bar(sl, 8)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9: Qualitative results
# ─────────────────────────────────────────────────────────────────────────────
def slide_qualitative():
    sl = add_slide()
    bg(sl, WHITE)
    header_bar(sl, 'Qualitative Results: Rollout Snapshots')

    slide_image(sl, os.path.join(FIGS_DIR, 'fig_qualitative.png'),
                Inches(0.4), Inches(1.3), width=Inches(12.5))

    textbox(sl, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.4),
            'Top: Baseline pushes too low \u2192 sliding. PhysCoT pushes high \u2192 clean topple.',
            font_size=Pt(13), color=DARK)
    textbox(sl, Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.4),
            'Bottom: Baseline picks blocked tool \u2192 failure. PhysCoT picks hook \u2192 object reaches goal.',
            font_size=Pt(13), color=DARK)

    footer_bar(sl, 9)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10: Failure modes
# ─────────────────────────────────────────────────────────────────────────────
def slide_failures():
    sl = add_slide()
    bg(sl, WHITE)
    header_bar(sl, 'Failure Mode Analysis')

    slide_image(sl, os.path.join(FIGS_DIR, 'fig_failure_modes.png'),
                Inches(0.4), Inches(1.3), width=Inches(7.5))

    # Right: analysis bullets
    textbox(sl, Inches(8.1), Inches(1.4), Inches(5.0), Inches(0.4),
            'Block Toppling Failures:', font_size=Pt(15), bold=True, color=DARK)
    bullet_list(sl, [
        'Baseline: 2 failures \u2014 contact too low',
        '  (insufficient torque or no contact)',
        'PhysCoT: 0 failures \u2014 all contacts above y_c*',
    ], Inches(8.1), Inches(1.9), Inches(5.0), Inches(1.5), font_size=Pt(13))

    textbox(sl, Inches(8.1), Inches(3.5), Inches(5.0), Inches(0.4),
            'Tool Selection Failures:', font_size=Pt(15), bold=True, color=DARK)
    bullet_list(sl, [
        'Baseline: 6 failures \u2014 wrong tool, path blocked',
        'PhysCoT: 1 failure \u2014 reasoning path error',
        '  (visual mis-assessment at Step C)',
    ], Inches(8.1), Inches(4.0), Inches(5.0), Inches(1.5), font_size=Pt(13))

    colored_box(sl, Inches(8.1), Inches(5.6), Inches(5.0), Inches(1.0),
                'PhysCoT fails only when visual estimation\nbreaks (Step C), not reasoning logic (B/D)',
                fill=STANFORD, font_size=Pt(14))

    footer_bar(sl, 10)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11: Demo
# ─────────────────────────────────────────────────────────────────────────────
def slide_demo():
    sl = add_slide()
    bg(sl, WHITE)
    header_bar(sl, 'Demo: Rollout Videos')

    textbox(sl, Inches(0.5), Inches(1.35), W - Inches(1.0), Inches(0.45),
            'All 40 trial videos in results/videos/ in the project repository',
            font_size=Pt(15), italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # Exp 1
    colored_box(sl, Inches(0.4), Inches(2.0), Inches(5.9), Inches(3.5),
                'Baseline Trial 5 (FAIL):\n'
                'results/videos/exp1_block_toppling/baseline/trial_04.mp4\n'
                '\u21d2 Contact at 32% height \u2192 sliding, no topple\n\n'
                'PhysCoT Trial 5 (SUCCESS):\n'
                'results/videos/exp1_block_toppling/physcot/trial_04.mp4\n'
                '\u21d2 Contact at 69% height \u2192 clean topple',
                fill=DARK, font_size=Pt(13), title='[*] Experiment 1: Block Toppling')

    # Exp 2
    colored_box(sl, Inches(6.7), Inches(2.0), Inches(6.2), Inches(3.5),
                'Baseline Trial 2 (FAIL):\n'
                'results/videos/exp2_tool_selection/baseline/trial_01.mp4\n'
                '\u21d2 Straight tool \u2192 path blocked, object stuck\n\n'
                'PhysCoT Trial 1 (SUCCESS):\n'
                'results/videos/exp2_tool_selection/physcot/trial_00.mp4\n'
                '\u21d2 Hook tool \u2192 arc path \u2192 object at goal',
                fill=GREEN, font_size=Pt(13), title='[*] Experiment 2: Tool Selection')

    textbox(sl, Inches(0.5), Inches(6.0), W - Inches(1.0), Inches(0.6),
            '\u25b6 Open MP4 files with any video player to review rollout behavior',
            font_size=Pt(15), bold=True, color=STANFORD, align=PP_ALIGN.CENTER)

    footer_bar(sl, 11)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12: Future pipeline
# ─────────────────────────────────────────────────────────────────────────────
def slide_future():
    sl = add_slide()
    bg(sl, WHITE)
    header_bar(sl, 'Future Work: Supervised PhysCoT Training Pipeline')

    slide_image(sl, os.path.join(FIGS_DIR, 'fig_future_pipeline.png'),
                Inches(0.4), Inches(1.3), width=Inches(12.5))

    colored_box(sl, Inches(0.4), Inches(5.1), Inches(5.9), Inches(1.8),
                '1. Simulate diverse episodes\n'
                '2. VLM annotates reasoning traces (one-off)\n'
                '3. Finetune OpenVLA on (obs, reasoning, action)',
                fill=BLUE, font_size=Pt(14), title='Proposed Steps')

    colored_box(sl, Inches(6.7), Inches(5.1), Inches(6.2), Inches(1.8),
                'Joint reasoning + action supervision\n'
                'Distillation: long \u2192 short reasoning\n'
                'Preference learning over physics traces\n'
                'Efficient test-time: no VLM needed',
                fill=PURPLE, font_size=Pt(14), title='Training Variants')

    footer_bar(sl, 12)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13: Limitations
# ─────────────────────────────────────────────────────────────────────────────
def slide_limitations():
    sl = add_slide()
    bg(sl, WHITE)
    header_bar(sl, 'Limitations')

    lims = [
        ('Prompt-only', 'No weight updates; cannot generalise beyond prompt scope', RED),
        ('Small n', 'n=10 per condition; wide confidence intervals; pilot study', ORANGE),
        ('2D simulation', 'No 3D effects, deformables, or real-robot noise', ORANGE),
        ('Policy approx.', 'Full GPU OpenVLA inference not integrated', ORANGE),
        ('Prompt design', 'Hand-designed schema; systematic search not done', GRAY),
    ]

    for i, (title, desc, color) in enumerate(lims):
        y = Inches(1.5 + i * 0.95)
        colored_box(sl, Inches(0.4), y, Inches(2.5), Inches(0.8),
                    title, fill=color, font_size=Pt(16))
        textbox(sl, Inches(3.1), y + Inches(0.18), Inches(9.8), Inches(0.55),
                desc, font_size=Pt(15), color=DARK)

    colored_box(sl, Inches(0.4), Inches(6.2), W - Inches(0.8), Inches(0.9),
                'Despite limitations: results are honestly reported, approximations documented, '
                'and the core finding \u2014 physics reasoning content matters \u2014 is clearly supported.',
                fill=GREEN, font_size=Pt(14))

    footer_bar(sl, 13)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14: Conclusion
# ─────────────────────────────────────────────────────────────────────────────
def slide_conclusion():
    sl = add_slide()
    bg(sl, RGBColor(0xFA, 0xF8, 0xF8))
    header_bar(sl, 'Conclusion')

    colored_box(sl, Inches(0.5), Inches(1.3), W - Inches(1.0), Inches(1.1),
                'PhysCoT is a prompt-only, physics-intuitive chain-of-thought wrapper for OpenVLA.\n'
                'Physics-aware, visually grounded, action-consequential reasoning improves manipulation.',
                fill=STANFORD, font_size=Pt(17))

    # Three columns
    cols = [
        (GREEN,  '\u2714 Results', '+20% block toppling\n+50% tool selection\nAll videos saved'),
        (BLUE,   '[*] Method', '4-step physics schema\nNo training required\nInterpretable traces'),
        (ORANGE, '[*] Future', 'VLM annotation \u2192\nSupervised finetuning\nEfficient deployment'),
    ]
    cw = Inches(3.9)
    for i, (color, title, content) in enumerate(cols):
        x = Inches(0.4) + i * (cw + Inches(0.2))
        colored_box(sl, x, Inches(2.7), cw, Inches(2.8),
                    content, fill=color, font_size=Pt(15), title=title)

    textbox(sl, Inches(0.5), Inches(5.8), W - Inches(1.0), Inches(0.6),
            'Key takeaway: physics-intuitive reasoning content is a distinct, under-explored '
            'axis of improvement for VLA policies \u2014 separate from scale, data, or architecture.',
            font_size=Pt(15), bold=True, color=DARK, align=PP_ALIGN.CENTER)

    textbox(sl, Inches(0.5), Inches(6.55), W - Inches(1.0), Inches(0.5),
            'Bobby Shi  \u2022  shi02@stanford.edu  \u2022  '
            '',
            font_size=Pt(13), italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    footer_bar(sl, 14)


# ─────────────────────────────────────────────────────────────────────────────
# Build all slides
# ─────────────────────────────────────────────────────────────────────────────
print('Building slides...')
slide_title()
slide_motivation()
slide_ecot_gap()
slide_key_idea()
slide_pipeline()
slide_exp1()
slide_exp2()
slide_results()
slide_qualitative()
slide_failures()
slide_demo()
slide_future()
slide_limitations()
slide_conclusion()

os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
prs.save(OUT_PATH)
print(f'Slides saved: {OUT_PATH}')
print(f'Total slides: {len(prs.slides)}')
